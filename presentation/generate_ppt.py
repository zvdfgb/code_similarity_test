from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
import os

def set_font(target, size, is_bold=False, color=None, font_name='Microsoft YaHei'):
    """
    Sets font properties for a Paragraph or a Run.
    If target is a Paragraph, it applies to all existing runs and the paragraph default.
    """
    def _apply_style(font_obj, element):
        font_obj.name = font_name
        font_obj.size = Pt(size)
        font_obj.bold = is_bold
        if color:
            font_obj.color.rgb = color
        
        # Set East Asian and Latin fonts explicitly to avoid garbled text
        try:
            rPr = element.get_or_add_rPr()
            rPr.rFonts.set(qn('a:ea'), font_name)
            rPr.rFonts.set(qn('a:latin'), font_name)
        except Exception:
            pass

    if hasattr(target, 'runs'): # Target is a Paragraph
        # 1. Apply to paragraph default properties
        # Note: target.font access on Paragraph gives the default font for the paragraph
        _apply_style(target.font, target._pPr)
        
        # 2. Apply to all existing runs in the paragraph
        for run in target.runs:
            _apply_style(run.font, run._r)
            
    elif hasattr(target, 'font'): # Target is a Run
        _apply_style(target.font, target._r)

def add_slide_title(slide, text):
    shapes = slide.shapes
    # Title strip
    shape = shapes.add_shape(1, Inches(0), Inches(0.4), Inches(0.2), Inches(0.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0, 51, 102) # Dark Blue
    shape.line.fill.background()
    
    # Title text
    # Increased width to prevent wrapping
    tb = shapes.add_textbox(Inches(0.4), Inches(0.3), Inches(9.5), Inches(1))
    p = tb.text_frame.paragraphs[0]
    p.text = text
    set_font(p, 28, True, RGBColor(0, 51, 102))

def add_footer(slide, text="C语言代码相似度检测系统", number=None):
    shapes = slide.shapes
    # Bottom line
    shape = shapes.add_shape(1, Inches(0.5), Inches(7), Inches(9), Inches(0.02))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(200, 200, 200)
    shape.line.fill.background()
    
    # Footer text
    tb = shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(5), Inches(0.4))
    p = tb.text_frame.paragraphs[0]
    p.text = text
    set_font(p, 10, False, RGBColor(150, 150, 150))
    
    if number:
        tb_num = shapes.add_textbox(Inches(9), Inches(7.1), Inches(0.5), Inches(0.4))
        p = tb_num.text_frame.paragraphs[0]
        p.text = str(number)
        p.alignment = PP_ALIGN.RIGHT
        set_font(p, 10, False, RGBColor(150, 150, 150))

def create_presentation():
    prs = Presentation()
    # Standard 4:3 is default, let's ensure it's Widescreen 16:9 for modern layout if needed
    # But python-pptx default is 4:3 (10x7.5 inches). We will stick to 10x7.5 layout logic used below.
    
    # --- 1. Title Slide ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Left dark panel
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(3.5), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0, 51, 102)
    shape.line.fill.background()
    
    # Title (Right side)
    tb = slide.shapes.add_textbox(Inches(4), Inches(2.5), Inches(5.5), Inches(2))
    p = tb.text_frame.paragraphs[0]
    p.text = "C语言代码\n相似度检测系统"
    p.line_spacing = 1.2
    set_font(p, 36, True, RGBColor(50, 50, 50))
    
    # Subtitle
    p = tb.text_frame.add_paragraph()
    p.text = "基于向量空间模型的查重算法实现"
    p.space_before = Pt(12)
    set_font(p, 18, False, RGBColor(100, 100, 100))
    
    # Info
    tb_info = slide.shapes.add_textbox(Inches(4), Inches(6), Inches(5.5), Inches(1))
    p = tb_info.text_frame.paragraphs[0]
    p.text = "2025年12月19日 | 课程设计答辩"
    set_font(p, 14, False, RGBColor(120, 120, 120))

    # --- 2. Background & Motivation ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "项目背景与意义")
    add_footer(slide, number=1)
    
    content = [
        ("当前痛点", "高校作业抄袭、工程代码重复利用率低，传统的人工比对费时费力。"),
        ("现有挑战", "简单的文本比对（Diff工具）无法处理格式化、重命名变量、调整函数顺序等修改手段。"),
        ("项目目标", "实现一个自动化、高鲁棒性的检测工具，关注代码的“逻辑骨架”而非“文本皮囊”。")
    ]
    
    top = 1.6
    for title, desc in content:
        # Icon/Box indicator
        s = slide.shapes.add_shape(1, Inches(0.8), Inches(top), Inches(0.15), Inches(0.6))
        s.fill.solid()
        s.fill.fore_color.rgb = RGBColor(0, 102, 204)
        s.line.fill.background()
        
        # Title
        tb = slide.shapes.add_textbox(Inches(1.1), Inches(top - 0.1), Inches(8.5), Inches(0.5))
        p = tb.text_frame.paragraphs[0]
        p.text = title
        set_font(p, 16, True, RGBColor(0, 0, 0))
        
        # Desc
        tb2 = slide.shapes.add_textbox(Inches(1.1), Inches(top + 0.35), Inches(8.5), Inches(0.8))
        tf = tb2.text_frame
        tf.word_wrap = True
        p2 = tf.paragraphs[0]
        p2.text = desc
        set_font(p2, 14, False, RGBColor(80, 80, 80))
        
        top += 1.6

    # --- 3. System Architecture ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "系统整体架构")
    add_footer(slide, number=2)
    
    if os.path.exists('flowchart.png'):
        # Adjusted size and position to fit better
        slide.shapes.add_picture('flowchart.png', Inches(0.5), Inches(1.5), width=Inches(6))
    
    # Side notes (Right side)
    tb = slide.shapes.add_textbox(Inches(6.8), Inches(2), Inches(3), Inches(4))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "流程解析"
    set_font(p, 16, True, RGBColor(0, 51, 102))
    
    bullets = ["源代码输入", "噪声清洗", "特征提取", "空间映射", "距离计算"]
    for b in bullets:
        p = tf.add_paragraph()
        p.text = f"• {b}"
        p.space_before = Pt(12)
        set_font(p, 14, False, RGBColor(60, 60, 60))

    # --- 4. Technical Detail: Preprocessing ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "关键技术一：预处理")
    add_footer(slide, number=3)
    
    # Concept
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.5), Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    p.text = "目的：消除代码风格差异，还原代码最纯粹的逻辑形态。"
    set_font(p, 16, False, RGBColor(50, 50, 50))
    
    # Comparison Visual
    # Left: Code Before
    box_l = slide.shapes.add_shape(1, Inches(0.8), Inches(2.2), Inches(3.8), Inches(3))
    box_l.fill.solid()
    box_l.fill.fore_color.rgb = RGBColor(245, 245, 245)
    box_l.line.color.rgb = RGBColor(200, 200, 200)
    
    tf_l = box_l.text_frame
    tf_l.margin_left = Inches(0.1)
    tf_l.margin_top = Inches(0.1)
    p = tf_l.paragraphs[0]
    p.text = "// Calculate sum\nint main() {\n  int a = 10; /* Init */\n  return a + 5;\n}"
    set_font(p, 12, False, RGBColor(0, 0, 0), 'Consolas')
    
    # Arrow
    arrow = slide.shapes.add_shape(33, Inches(4.8), Inches(3.5), Inches(0.6), Inches(0.4)) # Right Arrow
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = RGBColor(0, 102, 204)
    arrow.line.fill.background()
    
    # Right: Code After
    box_r = slide.shapes.add_shape(1, Inches(5.6), Inches(2.2), Inches(3.8), Inches(3))
    box_r.fill.solid()
    box_r.fill.fore_color.rgb = RGBColor(225, 240, 255)
    box_r.line.color.rgb = RGBColor(0, 102, 204)
    
    tf_r = box_r.text_frame
    tf_r.margin_left = Inches(0.1)
    tf_r.margin_top = Inches(0.1)
    tf_r.word_wrap = True
    p = tf_r.paragraphs[0]
    p.text = "int main ( ) { int a = ; return a + ; }"
    set_font(p, 12, False, RGBColor(0, 0, 0), 'Consolas')
    
    # Explanation
    tb_ex = slide.shapes.add_textbox(Inches(0.8), Inches(5.5), Inches(8.5), Inches(1))
    p = tb_ex.text_frame.paragraphs[0]
    p.text = "处理动作：\n1. 移除所有注释\n2. 压缩空白字符\n3. 移除字符串/数字常量"
    set_font(p, 14, True, RGBColor(80, 80, 80))

    # --- 5. Technical Detail: Vectorization ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "关键技术二：特征向量化")
    add_footer(slide, number=4)
    
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.5), Inches(0.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "策略：采用词袋模型 (Bag of Words)，忽略变量名，关注保留字与运算符。"
    set_font(p, 16, False, RGBColor(50, 50, 50))
    
    # 35-Dim Feature list visual
    s = slide.shapes.add_shape(1, Inches(0.8), Inches(2.3), Inches(8.5), Inches(1.5))
    s.fill.solid()
    s.fill.fore_color.rgb = RGBColor(255, 250, 240)
    s.line.dash_style = 1 # Solid
    s.line.color.rgb = RGBColor(200, 150, 100)
    
    tf = s.text_frame
    tf.margin_left = Inches(0.2)
    p = tf.paragraphs[0]
    p.text = "特征向量维度 (35维):"
    set_font(p, 14, True, RGBColor(200, 100, 0))
    
    p2 = tf.add_paragraph()
    p2.text = "[int, char, if, else, while, for, return, +, -, *, /, ==, !=, &&, ||, ...]"
    set_font(p2, 12, False, RGBColor(0, 0, 0), 'Consolas')
    
    # Why ignore variables?
    tb_why = slide.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(8.5), Inches(2.5))
    tf = tb_why.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "💡 创新点：为什么要忽略用户自定义变量？"
    set_font(p, 16, True, RGBColor(0, 51, 102))
    
    p2 = tf.add_paragraph()
    p2.text = "原因：如果两个程序逻辑不同（如冒泡排序 vs 快速排序），但定义了相同数量的变量（i, j, temp），会造成虚假的高相似度。"
    p2.space_before = Pt(10)
    set_font(p2, 14, False, RGBColor(60, 60, 60))
    
    p3 = tf.add_paragraph()
    p3.text = "效果：忽略变量名后，算法更聚焦于 if/while/for 等控制流结构，准确率显著提升。"
    p3.space_before = Pt(6)
    set_font(p3, 14, False, RGBColor(60, 60, 60))

    # --- 6. Core Algorithm ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "核心算法：余弦相似度")
    add_footer(slide, number=5)
    
    # Formula Box
    shape = slide.shapes.add_shape(1, Inches(1.5), Inches(2.0), Inches(7), Inches(1.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    shape.line.color.rgb = RGBColor(0, 51, 102)
    shape.line.width = Pt(2)
    
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = "Similarity = (A · B) / (||A|| × ||B||)"
    p.alignment = PP_ALIGN.CENTER
    set_font(p, 20, True, RGBColor(0, 0, 0))

    # Geometric Interpretation
    tb = slide.shapes.add_textbox(Inches(1.5), Inches(4.0), Inches(7), Inches(2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "几何意义："
    set_font(p, 16, True, RGBColor(0, 51, 102))
    
    p2 = tf.add_paragraph()
    p2.text = "将代码映射为多维空间中的向量，计算两个向量之间的夹角余弦值。"
    p2.space_before = Pt(10)
    set_font(p2, 14, False, RGBColor(60, 60, 60))
    
    p3 = tf.add_paragraph()
    p3.text = "• 夹角为0度 -> 余弦值为1 -> 完全相似"
    p3.space_before = Pt(5)
    set_font(p3, 14, False, RGBColor(60, 60, 60))
    
    p4 = tf.add_paragraph()
    p4.text = "• 夹角为90度 -> 余弦值为0 -> 完全不相关"
    p4.space_before = Pt(5)
    set_font(p4, 14, False, RGBColor(60, 60, 60))

    # --- 7. Evaluation: Test Design ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "实验设计与验证")
    add_footer(slide, number=6)
    
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.5), Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    p.text = "为了验证鲁棒性，设计了三组对照实验："
    set_font(p, 16, False, RGBColor(50, 50, 50))
    
    groups = [
        ("Group 1: 仅修改注释", "逻辑完全一致，仅将英文注释改为中文，或删除注释。", "预期：100%"),
        ("Group 2: 函数乱序", "保持所有函数内容不变，仅打乱函数在文件中的定义顺序。", "预期：100%"),
        ("Group 3: 不同算法", "完全不同的任务（如：冒泡排序 vs 斐波那契数列）。", "预期：< 60%")
    ]
    
    top = 2.2
    for title, desc, expect in groups:
        # Box container
        s = slide.shapes.add_shape(1, Inches(0.8), Inches(top), Inches(8.5), Inches(1.3))
        s.fill.solid()
        s.fill.fore_color.rgb = RGBColor(248, 248, 248)
        s.line.color.rgb = RGBColor(220, 220, 220)
        
        tf = s.text_frame
        tf.margin_top = Inches(0.1)
        tf.margin_left = Inches(0.2)
        
        p = tf.paragraphs[0]
        p.text = title
        set_font(p, 14, True, RGBColor(0, 51, 102))
        
        p2 = tf.add_paragraph()
        p2.text = desc
        set_font(p2, 12, False, RGBColor(80, 80, 80))
        
        p3 = tf.add_paragraph()
        p3.text = f"目标结果: {expect}"
        p3.space_before = Pt(3)
        set_font(p3, 12, True, RGBColor(200, 100, 0))
        
        top += 1.5

    # --- 8. Evaluation: Results ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "实验结果分析")
    add_footer(slide, number=7)
    
    # Table
    rows = 4
    cols = 4
    table_shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.8), Inches(9), Inches(2.2))
    table = table_shape.table
    
    headers = ["测试组", "对比文件", "相似度得分", "判定结论"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 51, 102)
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.alignment = PP_ALIGN.CENTER
        set_font(p, 12, True, RGBColor(255, 255, 255))
        
    data = [
        ("Group 1", "test1.c / test2.c", "1.0000", "极高相似 (抄袭)"),
        ("Group 2", "test3.c / test4.c", "1.0000", "极高相似 (抄袭)"),
        ("Group 3", "test5.c / test6.c", "0.5824", "低相似 (安全)")
    ]
    
    for row_idx, row_data in enumerate(data):
        for col_idx, item in enumerate(row_data):
            cell = table.cell(row_idx+1, col_idx)
            p = cell.text_frame.paragraphs[0]
            p.text = item
            p.alignment = PP_ALIGN.CENTER
            set_font(p, 12, False, RGBColor(0, 0, 0))

    # Analysis Text
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(2))
    p = tb.text_frame.paragraphs[0]
    p.text = "数据解读："
    set_font(p, 14, True, RGBColor(0, 102, 204))
    
    bullets = [
        "Group 1 (1.0000): 系统成功过滤了所有注释，验证了预处理模块的有效性。",
        "Group 2 (1.0000): 词袋模型忽略了代码顺序，因此简单的函数重排无法欺骗检测系统。",
        "Group 3 (0.5824): 对于完全不同的算法，即使共享基础关键字，得分依然显著偏低，区分度良好。"
    ]
    for b in bullets:
        p = tb.text_frame.add_paragraph()
        p.text = f"• {b}"
        p.space_before = Pt(8)
        set_font(p, 12, False, RGBColor(60, 60, 60))

    # --- 9. Implementation ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "工程实现与演示")
    add_footer(slide, number=8)
    
    # Left: Tech Stack
    s = slide.shapes.add_shape(1, Inches(0.5), Inches(1.8), Inches(4), Inches(2.5))
    s.fill.solid()
    s.fill.fore_color.rgb = RGBColor(240, 248, 255)
    s.line.color.rgb = RGBColor(200, 220, 240)
    
    tf = s.text_frame
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.2)
    p = tf.paragraphs[0]
    p.text = "技术栈"
    set_font(p, 14, True, RGBColor(0, 51, 102))
    
    stack = ["开发语言: C (Standard C11)", "编译工具: GCC / Makefile", "内存管理: 手动 (malloc/free)", "项目架构: 模块化 (头文件分离)"]
    for item in stack:
        p = tf.add_paragraph()
        p.text = f"• {item}"
        p.space_before = Pt(8)
        set_font(p, 12, False, RGBColor(0, 0, 0))
        
    # Right: Terminal Demo
    s2 = slide.shapes.add_shape(1, Inches(4.8), Inches(1.8), Inches(4.7), Inches(2.5))
    s2.fill.solid()
    s2.fill.fore_color.rgb = RGBColor(40, 44, 52) # Dark terminal
    
    tf = s2.text_frame
    tf.margin_left = Inches(0.1)
    tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = "> ./code_checker test1.c test2.c"
    set_font(p, 12, True, RGBColor(0, 255, 0), 'Consolas')
    p = tf.add_paragraph()
    p.text = "\n[INFO] Reading files...\n[INFO] Preprocessing...\n[INFO] Vectorizing...\n------------------------------\n[RESULT] Similarity: 1.0000\n[CONCLUSION] High Resemblance."
    set_font(p, 10, False, RGBColor(200, 200, 200), 'Consolas')

    # --- 10. Conclusion ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "总结与展望")
    add_footer(slide, number=9)
    
    # Summary Box
    tb_l = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.5), Inches(2))
    p = tb_l.text_frame.paragraphs[0]
    p.text = "项目总结"
    set_font(p, 18, True, RGBColor(0, 51, 102))
    
    p = tb_l.text_frame.add_paragraph()
    p.text = "本项目成功实现了一个基于特征向量的高效代码查重工具。相比传统文本比对，它更能抵抗格式化、注释修改等干扰手段，具有较高的鲁棒性和实用价值。"
    p.space_before = Pt(12)
    set_font(p, 14, False, RGBColor(60, 60, 60))
    
    # Future Box
    tb_r = slide.shapes.add_textbox(Inches(0.8), Inches(4.0), Inches(8.5), Inches(2.5))
    p = tb_r.text_frame.paragraphs[0]
    p.text = "未来改进方向"
    set_font(p, 18, True, RGBColor(200, 100, 0))
    
    futures = [
        "引入 抽象语法树 (AST): 检测更复杂的逻辑修改（如 while 转 for）。",
        "变量依赖分析: 防止恶意的变量名整体替换。",
        "多语言支持: 扩展对 Python, Java 等语言的解析支持。"
    ]
    for f in futures:
        p = tb_r.text_frame.add_paragraph()
        p.text = f"• {f}"
        p.space_before = Pt(10)
        set_font(p, 14, False, RGBColor(60, 60, 60))

    # --- 11. End Slide ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Center Text
    tb = slide.shapes.add_textbox(Inches(0), Inches(3), Inches(10), Inches(1.5))
    p = tb.text_frame.paragraphs[0]
    p.text = "感谢聆听\n请老师批评指正"
    p.alignment = PP_ALIGN.CENTER
    set_font(p, 32, True, RGBColor(0, 51, 102))

    prs.save('Project_Presentation.pptx')
    print("Final refined presentation generated.")

if __name__ == "__main__":
    create_presentation()
